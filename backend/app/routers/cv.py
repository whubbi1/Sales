from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import text
from app.database import get_db
from io import BytesIO
import uuid, json

router = APIRouter()

async def _table_exists(db: AsyncSession, name: str) -> bool:
    r = await db.execute(text("SELECT to_regclass(:n)"), {"n": name})
    return r.scalar() is not None

@router.get("")
async def list_all_cvs(db: AsyncSession = Depends(get_db)):
    from app.routers.settings import list_users as _list_users
    users_resp = await _list_users(db)
    users = users_resp.get("users", [])

    cr = await db.execute(text("SELECT email, title FROM employee_cv"))
    cv_titles = {row.email: row.title for row in cr.fetchall()}
    er = await db.execute(text("SELECT user_email, COUNT(*) AS c FROM employee_cv_experience GROUP BY user_email"))
    exp_counts = {row.user_email: row.c for row in er.fetchall()}

    result = []
    for u in users:
        email = u.get("email")
        result.append({
            **u,
            "cv_title": cv_titles.get(email) or u.get("job_title") or "",
            "has_cv": email in cv_titles,
            "experience_count": exp_counts.get(email, 0),
        })
    return {"users": result}

@router.get("/{email}")
async def get_cv(email: str, db: AsyncSession = Depends(get_db)):
    r = await db.execute(text("SELECT * FROM employee_cv WHERE email = :email"), {"email": email})
    row = r.fetchone()
    if row:
        cv = dict(row._mapping)
        cv["id"] = str(cv["id"])
    else:
        # Pre-fill from user_profiles if this employee hasn't set up a CV yet
        pr = await db.execute(text("SELECT first_name, last_name, job_title FROM user_profiles WHERE email = :email"), {"email": email})
        p = pr.fetchone()
        cv = {
            "id": None,
            "email": email,
            "first_name": p.first_name if p else "",
            "last_name": p.last_name if p else "",
            "title": p.job_title if p else "",
            "short_description": "",
            "skills": [],
            "languages": [],
        }

    er = await db.execute(text("""
        SELECT * FROM employee_cv_experience
        WHERE user_email = :email
        ORDER BY sort_order ASC, start_date DESC
    """), {"email": email})
    experiences = []
    for exp in er.fetchall():
        d = dict(exp._mapping)
        d["id"] = str(d["id"])
        experiences.append(d)
    cv["experiences"] = experiences

    # Read-only reflections from Training / Certifications (Phase 3). Tables may not exist yet.
    trainings, certifications = [], []
    if await _table_exists(db, "trainings"):
        tr = await db.execute(text("SELECT * FROM trainings WHERE user_email = :email ORDER BY training_date DESC"), {"email": email})
        trainings = [dict(t._mapping) for t in tr.fetchall()]
    if await _table_exists(db, "certifications"):
        cr = await db.execute(text("SELECT * FROM certifications WHERE user_email = :email ORDER BY cert_date DESC"), {"email": email})
        certifications = [dict(c._mapping) for c in cr.fetchall()]
    cv["trainings"] = trainings
    cv["certifications"] = certifications

    return {"cv": cv}

@router.put("/{email}")
async def update_cv(email: str, data: dict, db: AsyncSession = Depends(get_db)):
    await db.execute(text("""
        INSERT INTO employee_cv (id, email, first_name, last_name, title, short_description, skills, languages, created_at, updated_at)
        VALUES (gen_random_uuid(), :email, :first_name, :last_name, :title, :short_description, CAST(:skills AS JSON), CAST(:languages AS JSON), NOW(), NOW())
        ON CONFLICT (email) DO UPDATE SET
            first_name = EXCLUDED.first_name,
            last_name = EXCLUDED.last_name,
            title = EXCLUDED.title,
            short_description = EXCLUDED.short_description,
            skills = EXCLUDED.skills,
            languages = EXCLUDED.languages,
            updated_at = NOW()
    """), {
        "email": email,
        "first_name": data.get("first_name", ""),
        "last_name": data.get("last_name", ""),
        "title": data.get("title", ""),
        "short_description": data.get("short_description", ""),
        "skills": json.dumps(data.get("skills") or []),
        "languages": json.dumps(data.get("languages") or []),
    })
    await db.commit()
    return {"status": "ok"}

@router.post("/{email}/experience")
async def create_experience(email: str, data: dict, db: AsyncSession = Depends(get_db)):
    exp_id = str(uuid.uuid4())
    await db.execute(text("""
        INSERT INTO employee_cv_experience
            (id, user_email, job_title, company, start_date, end_date, location, description, sort_order, created_at, updated_at)
        VALUES
            (CAST(:id AS UUID), :email, :job_title, :company, :start_date, :end_date, :location, :description, :sort_order, NOW(), NOW())
    """), {
        "id": exp_id,
        "email": email,
        "job_title": data.get("job_title", ""),
        "company": data.get("company", ""),
        "start_date": data.get("start_date", ""),
        "end_date": data.get("end_date", ""),
        "location": data.get("location", ""),
        "description": data.get("description", ""),
        "sort_order": data.get("sort_order", 0),
    })
    await db.commit()
    return {"status": "ok", "id": exp_id}

@router.put("/{email}/experience/{eid}")
async def update_experience(email: str, eid: str, data: dict, db: AsyncSession = Depends(get_db)):
    await db.execute(text("""
        UPDATE employee_cv_experience SET
            job_title = :job_title,
            company = :company,
            start_date = :start_date,
            end_date = :end_date,
            location = :location,
            description = :description,
            updated_at = NOW()
        WHERE id = CAST(:id AS UUID) AND user_email = :email
    """), {
        "id": eid,
        "email": email,
        "job_title": data.get("job_title", ""),
        "company": data.get("company", ""),
        "start_date": data.get("start_date", ""),
        "end_date": data.get("end_date", ""),
        "location": data.get("location", ""),
        "description": data.get("description", ""),
    })
    await db.commit()
    return {"status": "ok"}

@router.delete("/{email}/experience/{eid}")
async def delete_experience(email: str, eid: str, db: AsyncSession = Depends(get_db)):
    await db.execute(text("DELETE FROM employee_cv_experience WHERE id = CAST(:id AS UUID) AND user_email = :email"), {"id": eid, "email": email})
    await db.commit()
    return {"status": "ok"}

# ─── Export: Word (complete or curated) and PowerPoint (summary) ──────────────
async def _load_cv_for_export(email: str, db: AsyncSession) -> dict:
    resp = await get_cv(email, db)
    return resp["cv"]

def _cv_full_name(cv: dict) -> str:
    return f"{cv.get('first_name','')} {cv.get('last_name','')}".strip() or cv.get("email", "")

def _generate_cv_docx(cv: dict, experience_ids: list | None) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from app.routers.hr import _fetch_logo

    C_PRI = RGBColor(0x15, 0x60, 0x82)
    C_SEC = RGBColor(0x45, 0xB6, 0xE4)

    experiences = cv.get("experiences") or []
    if experience_ids:
        wanted = set(experience_ids)
        experiences = [e for e in experiences if e["id"] in wanted]

    doc = Document()
    for sec in doc.sections:
        sec.top_margin = Cm(2); sec.bottom_margin = Cm(2.5)
        sec.left_margin = Cm(2.5); sec.right_margin = Cm(2.5)
        header = sec.header
        hp = header.paragraphs[0]
        hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
        logo_bytes = _fetch_logo()
        if logo_bytes:
            run_logo = hp.add_run()
            run_logo.add_picture(BytesIO(logo_bytes), height=Cm(1.2))
        else:
            run_h = hp.add_run("WCOMPLY")
            run_h.font.bold = True; run_h.font.color.rgb = C_PRI; run_h.font.size = Pt(14)

    p = doc.add_paragraph()
    run = p.add_run(_cv_full_name(cv))
    run.font.size = Pt(22); run.font.bold = True; run.font.color.rgb = C_PRI
    p.paragraph_format.space_after = Pt(2)

    if cv.get("title"):
        pt = doc.add_paragraph()
        rt = pt.add_run(cv["title"])
        rt.font.size = Pt(13); rt.font.color.rgb = C_SEC; rt.font.bold = True
        pt.paragraph_format.space_after = Pt(8)

    def add_h(text_):
        ph = doc.add_paragraph()
        r = ph.add_run(text_)
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = C_SEC
        ph.paragraph_format.space_before = Pt(14); ph.paragraph_format.space_after = Pt(4)

    def add_body(text_):
        pb = doc.add_paragraph()
        pb.add_run(text_).font.size = Pt(10)
        pb.paragraph_format.space_after = Pt(4)

    def add_bullet(text_):
        pb = doc.add_paragraph(style='List Bullet')
        pb.add_run(text_).font.size = Pt(10)

    if cv.get("short_description"):
        add_body(cv["short_description"])

    if cv.get("skills"):
        add_h("Compétences")
        add_body(", ".join(cv["skills"]))

    if cv.get("languages"):
        add_h("Languages")
        add_body(", ".join(cv["languages"]))

    if experiences:
        add_h("Experience")
        for exp in experiences:
            pe = doc.add_paragraph()
            re_ = pe.add_run(f"{exp.get('job_title','')} — {exp.get('company','')}")
            re_.font.bold = True; re_.font.size = Pt(11); re_.font.color.rgb = C_PRI
            pe.paragraph_format.space_after = Pt(1)
            dates = f"{exp.get('start_date') or '—'} → {exp.get('end_date') or 'Present'}"
            loc = f" · {exp['location']}" if exp.get("location") else ""
            pd = doc.add_paragraph()
            rd = pd.add_run(dates + loc)
            rd.font.size = Pt(9); rd.font.italic = True
            pd.paragraph_format.space_after = Pt(3)
            if exp.get("description"):
                add_body(exp["description"])

    if cv.get("trainings"):
        add_h("Trainings")
        for t in cv["trainings"]:
            add_bullet(f"{t.get('name','')} ({t.get('training_date') or '—'})")

    if cv.get("certifications"):
        add_h("Certifications")
        for c in cv["certifications"]:
            add_bullet(f"{c.get('name','')} ({c.get('cert_date') or '—'})")

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()

def _generate_cv_pptx(cv: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from app.routers.hr import _fetch_logo

    C_PRI = RGBColor(0x15, 0x60, 0x82)
    C_SEC = RGBColor(0x45, 0xB6, 0xE4)
    C_TXT = RGBColor(0x3F, 0x3F, 0x3F)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — title
    s1 = prs.slides.add_slide(blank)
    band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3))
    band.fill.solid(); band.fill.fore_color.rgb = C_PRI
    band.line.fill.background()
    logo_bytes = _fetch_logo()
    if logo_bytes:
        s1.shapes.add_picture(BytesIO(logo_bytes), Inches(0.6), Inches(0.5), height=Inches(0.6))
    name_box = s1.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(1))
    tf = name_box.text_frame
    tf.text = _cv_full_name(cv)
    tf.paragraphs[0].font.size = Pt(40); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if cv.get("title"):
        title_box = s1.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12), Inches(0.6))
        tf2 = title_box.text_frame
        tf2.text = cv["title"]
        tf2.paragraphs[0].font.size = Pt(20); tf2.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if cv.get("short_description"):
        desc_box = s1.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(11.5), Inches(2))
        tf3 = desc_box.text_frame
        tf3.word_wrap = True
        tf3.text = cv["short_description"]
        tf3.paragraphs[0].font.size = Pt(14); tf3.paragraphs[0].font.color.rgb = C_TXT

    # Slide 2 — highlights
    s2 = prs.slides.add_slide(blank)
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_SEC
    bar.line.fill.background()
    heading = s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11), Inches(0.7))
    heading.text_frame.text = "Highlights"
    heading.text_frame.paragraphs[0].font.size = Pt(28); heading.text_frame.paragraphs[0].font.bold = True; heading.text_frame.paragraphs[0].font.color.rgb = C_PRI

    col_w = Inches(5.8)
    skills_box = s2.shapes.add_textbox(Inches(0.6), Inches(1.4), col_w, Inches(2.5))
    stf = skills_box.text_frame
    stf.word_wrap = True
    stf.text = "Compétences"
    stf.paragraphs[0].font.bold = True; stf.paragraphs[0].font.size = Pt(16); stf.paragraphs[0].font.color.rgb = C_SEC
    for s in (cv.get("skills") or []):
        p = stf.add_paragraph()
        p.text = f"• {s}"
        p.font.size = Pt(13); p.font.color.rgb = C_TXT

    lang_box = s2.shapes.add_textbox(Inches(6.8), Inches(1.4), col_w, Inches(2.5))
    ltf = lang_box.text_frame
    ltf.word_wrap = True
    ltf.text = "Languages"
    ltf.paragraphs[0].font.bold = True; ltf.paragraphs[0].font.size = Pt(16); ltf.paragraphs[0].font.color.rgb = C_SEC
    for l in (cv.get("languages") or []):
        p = ltf.add_paragraph()
        p.text = f"• {l}"
        p.font.size = Pt(13); p.font.color.rgb = C_TXT

    recent = (cv.get("experiences") or [])[:3]
    if recent:
        exp_box = s2.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(2.8))
        etf = exp_box.text_frame
        etf.word_wrap = True
        etf.text = "Recent Experience"
        etf.paragraphs[0].font.bold = True; etf.paragraphs[0].font.size = Pt(16); etf.paragraphs[0].font.color.rgb = C_SEC
        for exp in recent:
            p = etf.add_paragraph()
            p.text = f"{exp.get('job_title','')} — {exp.get('company','')}"
            p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = C_PRI

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

@router.get("/{email}/export/word")
async def export_word(email: str, experience_ids: str = None, db: AsyncSession = Depends(get_db)):
    cv = await _load_cv_for_export(email, db)
    ids = experience_ids.split(",") if experience_ids else None
    content = _generate_cv_docx(cv, ids)
    fname = f"{_cv_full_name(cv).replace(' ', '_') or 'CV'}.docx"
    return StreamingResponse(BytesIO(content),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'})

@router.get("/{email}/export/pptx")
async def export_pptx(email: str, db: AsyncSession = Depends(get_db)):
    cv = await _load_cv_for_export(email, db)
    content = _generate_cv_pptx(cv)
    fname = f"{_cv_full_name(cv).replace(' ', '_') or 'CV'}_Summary.pptx"
    return StreamingResponse(BytesIO(content),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'})
